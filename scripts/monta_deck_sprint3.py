# -*- coding: utf-8 -*-
"""Monta o deck da Sprint 3: HTML de cada slide, PNG de cada HTML e o .pptx final.

A ordem segue o template oficial da FIAP
(`assets/Templates/03Template_MVP_Preliminar_Challenge_2026_01_locaweb.pptx`), com os mesmos
títulos. Onde o conteúdo pede mais espaço, um item do template vira dois slides. Foi o que a
Sprint 2 fez, e tirou 5,00.

Quase tudo vem de arquivo pronto, e o builder só ordena e monta:

* **`prototipos/slides/mvp/abertura/`** traz a identificação, o contexto, o problema, a
  proposta, a gestão e a arquitetura, pelo mapa `ESCOLHIDOS`;
* **`prototipos/slides/mvp/deck/`** traz a análise e a modelagem, pela lista `ANALISE`, com
  os gráficos exportados dos notebooks;
* **`prototipos/slides/mvp/aplicacao/`** traz a divisória, as dez telas e o fecho, pelo mapa
  `ESCOLHIDOS_APP`. As telas são geradas por `scripts/monta_prints_app.py`;
* **este arquivo** escreve só o mockup.

O builder lê cada slide de onde ele está, sem copiar, para não existir uma segunda cópia
envelhecendo em silêncio.

Uso, com a aplicação já capturada por `captura_telas.py`, as telas montadas por
`monta_prints_app.py` e o quadro por `monta_quadro_trello.py`:

    .venv/Scripts/python scripts/monta_deck_sprint3.py

Duas regras atravessam o arquivo inteiro. **Todo número aqui saiu de célula executada**, e a
origem está anotada no comentário do slide. E métrica de avaliação de modelo é assunto de
slide, nunca de tela: o sistema simula um relógio parado em 01/10/2025 15h.
"""
from pathlib import Path

from deck_estilo import (
    ALTURA,
    LARGURA,
    pagina,
    slide,
)

RAIZ = Path(__file__).resolve().parents[1]
PRINTS = RAIZ / "sprints" / "sprint-3" / "prints"
SAIDA = RAIZ / "sprints" / "sprint-3" / "slides"
SAIDA_PNG = SAIDA / "png"
PPTX = RAIZ / "sprints" / "EC_Sprint_3_2TSCOA_Evidencias_Construcao_Cronos_SuperDataBros.pptx"
# Réplica do quadro Trello, gerada por `scripts/monta_quadro_trello.py`. Fica fora de
# `prints/` porque não é captura da aplicação, e o builder só a consome.
QUADRO = RAIZ / "sprints" / "sprint-3" / "quadro" / "png"
# Nenhuma URL do quadro aparece no deck: o quadro real ainda não está publicado, e simular
# uma barra de navegador com endereço que não existe seria evidência falsa.

# Os 38 slides de análise e modelagem, na ordem do `viewer.html`, escolhendo a primeira
# variação de cada um, que é a que o viewer abre.
ABERTURA = RAIZ / "prototipos" / "slides" / "mvp" / "abertura"
# A abertura foi refeita em três rodadas de protótipo e o Igor escolheu uma versão por
# slide. O builder lê o arquivo escolhido de onde ele está, em vez de reescrever o slide
# em Python: com as duas versões vivas, uma delas envelheceria em silêncio.
#
#   nome no deck  ->  arquivo em prototipos/slides/mvp/abertura/
ESCOLHIDOS = {
    "capa": "01-capa-B",
    "equipe": "02-equipe-B",
    "contexto": "03-contexto-E",
    "contexto-estouro": "03b-causas-H",
    "problema": "04-problema-D",
    "solucao": "05-solucao-D",
    "gestao-documentacao": "06-gestao-doc-D",
    "gestao-planejamento": "07-gestao-plano-D",
    "gestao-cartao": "07b-cartao-D",
    "fontes-de-dados": "08-arquitetura-D",
    "arquitetura-desenho": "09-desenho-D",
    "arquitetura-descricao": "10-descricao-A",
    "arquitetura-tecnologias": "11-tecnologias-A",
    # O template pede "entregar algoritmos, métodos, manipulações e transformações
    # utilizadas" (slide 13) e o enunciado pede "uso efetivo de código-fonte". O deck
    # mostrava o resultado do código em todo o bloco analítico e não mostrava o código.
    "codigo-fonte": "12-codigo-C",
}

DECK_ANALISE = RAIZ / "prototipos" / "slides" / "mvp" / "deck"
# Sufixo `-r5a` ou `-r5b` marca a versão escolhida pelo Igor na rodada 5, em que cada slide
# reprovado foi refeito em duas composições. Slide sem sufixo é o que ele não pediu para
# mexer. Quatro saíram do deck e continuam em disco: `d11a` e `d28a` a pedido dele, `d39a`
# porque o fechamento do bloco não estava comunicando, e `d31a` porque repetia o achado de
# familiaridade que já está no `d22a`, sobre outra base, o que produzia dois números para a
# mesma coisa. Ver `docs/auditoria-consistencia-deck.md`.
#
# `-r6b` é o único da rodada 6: o `d33a` voltou da rodada 5 com título de veredito e uma
# tabela de onze números repetindo o gráfico. A versão escolhida deixa a figura ocupar o
# slide até o pé, com uma linha de texto acima. A figura foi refeita por
# `scripts/figura_controles.py`, que confere cada valor contra o output do notebook 05.
ANALISE = [
    # Análise exploratória
    "d01a-r5a", "d02a-r5a", "d03m-r5b", "d04r-r5b", "d04m-r5a", "d05m", "d06m", "d07m",
    "d08m", "d09a", "d10a-r5a",
    # Previsão de volume
    "d12a", "d13a", "d14m", "d15m", "d15n-r5b", "d17a-r5b",
    # Risco de estouro de OLA
    "d20a", "d21a-r5a", "d22a-r5b", "d23a-r5b", "d24a-r5a", "d25a-r5a", "d26a-r5b",
    "d27a-r5a",
    # Causas e recorrência
    "d29a", "d30a-r5a", "d32a-r5b", "d33a-r6b",
    # Projeção do KPI e saúde por produto
    "d34a", "d35a-r5b", "d36a-r5b", "d37a-r5b", "d38a-r5a",
]

# O bloco da aplicação também vem de arquivo pronto. Os dez slides de tela são gerados por
# `scripts/monta_prints_app.py`, com os modais logo depois da aba de onde saem.
APLICACAO = RAIZ / "prototipos" / "slides" / "mvp" / "aplicacao"
TELAS_APP = [
    "tela-01-01-panorama", "tela-02-07-modal-briefing", "tela-03-02-previsao",
    "tela-04-03-projecao", "tela-05-10-modal-meta", "tela-06-04-fila",
    "tela-07-08-modal-escore", "tela-08-05-saude", "tela-09-09-modal-produto",
    "tela-10-06-causas",
]
ESCOLHIDOS_APP = {
    "divisoria-telas": "divisoria-r5a",
    "fecho": "fecho-r5b",
    **{n: n for n in TELAS_APP},
}

RODAPE_ESQ = "Cronos · Super Data Bros · 2TSCOA"
RODAPE_DIR = "Challenge FIAP 2026 com Locaweb"


def _img(nome: str) -> str:
    return (PRINTS / f"{nome}.png").resolve().as_uri()


def confere_quadro() -> None:
    """Os slides 06 e 07b da abertura carregam estas imagens por caminho relativo.

    O `<img>` quebrado não falha o build do Chromium, então a checagem é aqui: sem isso, o
    .pptx sairia com dois slides em branco e ninguém perceberia até abrir o arquivo.
    """
    for nome in ("quadro", "card-ppt"):
        caminho = QUADRO / f"{nome}.png"
        if not caminho.exists():
            raise FileNotFoundError(
                f"{caminho} não existe. Rode antes: "
                ".venv/Scripts/python scripts/monta_quadro_trello.py --png"
            )


def _sl(conteudo: str, *, tag: str, escuro: bool = False, classe: str = "") -> str:
    return slide(
        conteudo,
        tag=tag,
        rodape_esq=RODAPE_ESQ,
        rodape_dir=RODAPE_DIR,
        escuro=escuro,
        classe=classe,
    )


def s13_mockup() -> str:
    conteudo = (
        '<div style="text-align:center">'
        '<span class="eb" style="justify-content:center">A ferramenta em uso</span>'
        '<h1 class="tt" style="font-size:36px">O dia começa aqui</h1></div>'
        '<div class="nb"><div class="scr">'
        f'<img src="{_img("01-panorama")}" alt=""></div>'
        '<div class="base"></div></div>'
        '<p class="cap">A coordenação abre o Cronos, lê o briefing e já sabe onde agir. '
        "<b>Nenhuma pergunta foi feita.</b></p>"
    )
    return _sl(conteudo, tag="Mockup", classe="mk")


# ─────────────────────────────────────────────────────────────────────────────
def slides_proprios() -> dict[str, str]:
    """Os slides que nascem aqui. Abertura e bloco analítico vêm prontos de `prototipos/`."""
    return {"mockup": s13_mockup()}


def ordem() -> list[str]:
    """A sequência final do deck, na ordem do template da FIAP."""
    return (
        [
            "capa",                      # template 1 e 2
            "equipe",
            "contexto",                  # template 3: a operação e o prazo de OLA
            "contexto-estouro",          # template 3: por que o estouro é difícil de prever
            "problema",                  # template 4
            "solucao",                   # template 5
            "gestao-documentacao",       # template 6: a imagem do quadro, em tela cheia
            "gestao-planejamento",       # template 7: o cronograma das quatro entregas
            "gestao-cartao",             # template 7: uma atividade aberta dentro do quadro
            "fontes-de-dados",           # template 8
            "arquitetura-desenho",       # template 9
            "arquitetura-descricao",     # template 10
            "arquitetura-tecnologias",   # template 11
            "codigo-fonte",              # template 12 e 13: o código que produziu o resto
        ]
        + ANALISE                        # template 12 e 13: a evidência analítica
        + ["divisoria-telas"]            # template 13: a aplicação
        + TELAS_APP
        + ["mockup", "fecho"]            # template 14
    )


def escreve_html() -> list[tuple[str, Path]]:
    """Grava os slides próprios e devolve a sequência inteira resolvida em caminho.

    Slide do bloco analítico não é reescrito: ele é lido de onde está. Copiar o HTML para cá
    criaria uma segunda cópia que envelheceria em silêncio.
    """
    SAIDA.mkdir(parents=True, exist_ok=True)
    confere_quadro()
    proprios = slides_proprios()
    caminhos: list[tuple[str, Path]] = []
    for pos, nome in enumerate(ordem(), 1):
        if nome in proprios:
            destino = SAIDA / f"{pos:02d}-{nome}.html"
            destino.write_text(pagina(proprios[nome], nome), encoding="utf-8")
        elif nome in ESCOLHIDOS:
            destino = ABERTURA / f"{ESCOLHIDOS[nome]}.html"
            if not destino.exists():
                raise FileNotFoundError(f"slide da abertura não encontrado: {destino}")
        elif nome in ESCOLHIDOS_APP:
            destino = APLICACAO / f"{ESCOLHIDOS_APP[nome]}.html"
            if not destino.exists():
                raise FileNotFoundError(
                    f"slide da aplicação não encontrado: {destino}. Se for uma tela, rode "
                    "antes: .venv/Scripts/python scripts/monta_prints_app.py"
                )
        else:
            destino = DECK_ANALISE / f"{nome}.html"
            if not destino.exists():
                raise FileNotFoundError(f"slide do bloco analítico não encontrado: {destino}")
        caminhos.append((f"{pos:02d}-{nome}", destino))
    return caminhos


def renderiza_png(caminhos: list[tuple[str, Path]]) -> list[Path]:
    """Cada HTML vira um PNG em 2×, que é o que entra no .pptx."""
    from playwright.sync_api import sync_playwright

    SAIDA_PNG.mkdir(parents=True, exist_ok=True)
    pngs = []
    with sync_playwright() as pw:
        navegador = pw.chromium.launch(channel="chrome")
        page = navegador.new_context(
            viewport={"width": LARGURA, "height": ALTURA},
            device_scale_factor=2,
        ).new_page()
        for nome, html in caminhos:
            page.goto(html.resolve().as_uri(), wait_until="networkidle")
            page.evaluate("document.fonts.ready")
            page.wait_for_timeout(450)
            destino = SAIDA_PNG / f"{nome}.png"
            page.locator(".slide").first.screenshot(path=str(destino))
            pngs.append(destino)
        navegador.close()
    print(f"     {len(pngs)} slides renderizados")
    return pngs


def limpa_orfaos(caminhos: list[tuple[str, Path]], pngs: list[Path]) -> int:
    """Apaga HTML e PNG de numerações antigas.

    O nome de cada arquivo começa pela posição no deck, então inserir um slide no meio
    renumera tudo o que vem depois e deixa a versão anterior para trás. Sem esta limpeza a
    pasta acumula duas gerações do mesmo slide, e a mais velha é a que aparece primeiro
    quando alguém abre a pasta para conferir.
    """
    vivos = {c.name for _, c in caminhos} | {p.name for p in pngs}
    mortos = [
        p
        for p in list(SAIDA.glob("*.html")) + list(SAIDA_PNG.glob("*.png"))
        if p.name not in vivos
    ]
    for p in mortos:
        p.unlink()
    return len(mortos)


def monta_pptx(pngs: list[Path]) -> None:
    """Um slide por PNG, sangrado. É o mesmo formato entregue na Sprint 2."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(12192000)  # 13,333 in
    prs.slide_height = Emu(6858000)  # 7,5 in
    branco = prs.slide_layouts[6]
    for png in pngs:
        s = prs.slides.add_slide(branco)
        s.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    prs.save(PPTX)


def main() -> None:
    print("1/3 · escrevendo HTML dos slides próprios")
    caminhos = escreve_html()
    proprios = sum(1 for _, p in caminhos if p.parent == SAIDA)
    print(f"     {proprios} próprios + {len(caminhos)-proprios} do bloco analítico")
    print("2/3 · renderizando PNG")
    pngs = renderiza_png(caminhos)
    orfaos = limpa_orfaos(caminhos, pngs)
    if orfaos:
        print(f"     {orfaos} arquivos de numeração antiga removidos")
    print("3/3 · montando o .pptx")
    monta_pptx(pngs)
    print(f"Pronto: {PPTX.name} ({len(pngs)} slides)")


if __name__ == "__main__":
    main()
