"""
Extrai conteudo detalhado de um PPTX: texto, shapes, tabelas, imagens.
Salva markdown + diretorio de imagens extraidas.

Uso:
    python sprints/_extract_ppt_content.py <input.pptx> <output_dir>
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass


def shape_kind(shape):
    try:
        return shape.shape_type
    except Exception:
        return "?"


def extract(pptx_path: Path, out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)
    img_dir = out_dir / "images"
    img_dir.mkdir(exist_ok=True)

    prs = Presentation(str(pptx_path))
    md = out_dir / f"{pptx_path.stem}.md"

    lines = []
    w = prs.slide_width / 914400
    h = prs.slide_height / 914400
    lines.append(f"# {pptx_path.name}")
    lines.append("")
    lines.append(f"- Slide size: **{w:.2f} x {h:.2f} in**")
    lines.append(f"- Total slides: **{len(prs.slides)}**")
    lines.append("")

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"---")
        lines.append("")
        lines.append(f"## Slide {i} — layout: `{slide.slide_layout.name}`")
        lines.append("")

        for j, shape in enumerate(slide.shapes, 1):
            kind = shape_kind(shape)

            # texto
            if shape.has_text_frame:
                texts = []
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        texts.append(t)
                if texts:
                    full = " | ".join(texts)
                    name = shape.name or "?"
                    lines.append(f"**Text** (shape #{j} `{name}`):")
                    for t in texts:
                        lines.append(f"  > {t}")
                    lines.append("")

            # imagem
            if kind == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    ext = img.ext
                    out_img = img_dir / f"slide{i:02d}_shape{j:02d}.{ext}"
                    out_img.write_bytes(img.blob)
                    w_in = shape.width / 914400 if shape.width else 0
                    h_in = shape.height / 914400 if shape.height else 0
                    lines.append(f"**Picture** (shape #{j} `{shape.name}`): "
                                 f"{w_in:.2f}x{h_in:.2f}in -> `images/{out_img.name}`")
                    lines.append("")
                except Exception as e:
                    lines.append(f"**Picture** (shape #{j} `{shape.name}`): erro extraindo ({e})")
                    lines.append("")

            # tabela
            if shape.has_table:
                lines.append(f"**Table** (shape #{j} `{shape.name}`):")
                lines.append("")
                for row in shape.table.rows:
                    cells = [c.text.replace("\n", " ").strip()[:80] for c in row.cells]
                    lines.append(f"  | {' | '.join(cells)} |")
                lines.append("")

            # grafico
            try:
                if shape.has_chart:
                    lines.append(f"**Chart** (shape #{j}): tipo `{shape.chart.chart_type}`")
                    lines.append("")
            except Exception:
                pass

    md.write_text("\n".join(lines), encoding="utf-8")
    print(f"[OK] {pptx_path.name} -> {md.relative_to(out_dir.parent)} ({len(prs.slides)} slides)")


def main():
    if len(sys.argv) != 3:
        print("uso: python sprints/_extract_ppt_content.py <input.pptx> <output_dir>")
        sys.exit(1)
    pptx = Path(sys.argv[1])
    out = Path(sys.argv[2])
    extract(pptx, out)


if __name__ == "__main__":
    main()
