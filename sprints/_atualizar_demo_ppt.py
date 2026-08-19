"""
Substitui o slide 'Demo navegavel' (slide 24) pela versao LIGHT V3.
Encontra o slide pela posicao conhecida (24, indice 23) e troca a imagem.

Uso:
    python sprints/_atualizar_demo_ppt.py
"""
import sys
from pathlib import Path

from pptx import Presentation

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

SPRINTS = Path(__file__).parent
PROTO = SPRINTS.parent / "prototipos"
PPT_PATH = SPRINTS / "EC_Sprint_2_2TSCOA_arqsolucao_Cronos_SuperDataBros.pptx"
DEMO_NEW = PROTO / "gestao" / "screenshots" / "demo-light-v3.png"

SLIDE_DEMO_IDX = 23  # 0-based — slide 24 (Demo)


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        tag = sp.tag.split("}")[-1]
        if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
            spTree.remove(sp)


def main():
    prs = Presentation(str(PPT_PATH))
    total = len(prs.slides)
    print(f"  Slides totais: {total}")

    if SLIDE_DEMO_IDX >= total:
        print(f"  [ERR] slide #{SLIDE_DEMO_IDX + 1} fora do range")
        return

    slide = list(prs.slides)[SLIDE_DEMO_IDX]
    clear_slide(slide)
    slide.shapes.add_picture(
        str(DEMO_NEW),
        0, 0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    print(f"  Slide #{SLIDE_DEMO_IDX + 1} substituido por {DEMO_NEW.name}")

    prs.save(str(PPT_PATH))
    final_mb = PPT_PATH.stat().st_size / 1024 / 1024
    print(f"\n  PPTX salvo: {PPT_PATH.name}")
    print(f"  Tamanho: {final_mb:.1f} MB")


if __name__ == "__main__":
    main()
