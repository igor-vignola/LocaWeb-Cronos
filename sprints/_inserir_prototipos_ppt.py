"""
Insere os screenshots dos protótipos no PPTX da Sprint 2.

Estratégia robusta (sem remover slides):
  1. Os 5 placeholders existentes (slides 6-10) sao ESVAZIADOS e reaproveitados
     como divisores das 5 telas (Dashboard, Previsoes, Cascatas, Saude, Morning).
  2. Os 15 prints sao adicionados ao final do PPT.
  3. Reordena pra cada print ficar logo apos seu divisor.

Por que: o remove_slide do python-pptx gera "Duplicate name" no zip salvo —
esvaziar evita esse bug.

Uso:
    python sprints/_inserir_prototipos_ppt.py
"""
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

# ─── PATHS ───
SPRINTS = Path(__file__).parent
PROTO = SPRINTS.parent / "prototipos"
PPT_PATH = SPRINTS / "EC_Sprint_2_2TSCOA_arqsolucao_Cronos_SuperDataBros.pptx"
BACKUP = SPRINTS / "EC_Sprint_2_2TSCOA_arqsolucao_Cronos_SuperDataBros.bkp.pptx"
SHOTS = PROTO / "screenshots"

# ─── TELAS ───
# (nome_exibicao, prefixo, descricao, arquivo_divisor_screenshot)
TELAS = [
    ("Dashboard",        "dashboard",      "Visão geral · KPI mensal · cascatas ativas · tendência 30 dias",  "dashboard-v1.png"),
    ("Previsões",        "previsao",       "Forecast D+1 e D+7 · drivers SHAP · top produtos em risco",       "previsao-v1.png"),
    ("Cascatas",         "cascata",        "Acúmulo de P4/P5 antecedendo escaladas P3/P2",                    "cascata-v1.png"),
    ("Saúde · Produto",  "saude-produto",  "Score 0-100 por produto · ranking · tendência por janela",        "saude-produto-v1.png"),
    ("Morning brief",    "morning-brief",  "Resumo escrito automático · Ontem · Hoje · Ações sugeridas",      "morning-brief-v1.png"),
]

DIV_SHOTS = PROTO / "divisores" / "screenshots"

# ─── CORES ───
DARK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA3, 0xA3, 0xA3)
LIGHT_BG = RGBColor(0xEE, 0xF1, 0xF7)

# placeholders no PPT atual (slides 6-10 = indices 5-9)
PLACEHOLDER_IDXS = [5, 6, 7, 8, 9]
INSERT_AFTER = 4  # após slide #5 (indice 4)


# ─── HELPERS ───
def clear_slide(slide):
    """Remove TODAS as shapes de um slide (esvazia)."""
    spTree = slide.shapes._spTree
    # mantem apenas os shapes que sao placeholders padrao do layout? Não — mato tudo
    # alem dos elementos básicos (nvGrpSpPr, grpSpPr)
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    for sp in list(spTree):
        tag = sp.tag.split("}")[-1]
        if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
            spTree.remove(sp)


def add_solid_bg(slide, prs, color: RGBColor):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height,
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    # manda pra tras (logo apos nvGrpSpPr e grpSpPr)
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return rect


def add_text(slide, text, *, x, y, w, h, size, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Outfit"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def build_divisor(slide, prs, divisor_img: Path):
    """Transforma um slide existente (esvaziado) num divisor usando screenshot HTML em fullscreen."""
    clear_slide(slide)
    # imagem ocupa todo o slide
    slide.shapes.add_picture(
        str(divisor_img),
        0, 0,
        width=prs.slide_width,
        height=prs.slide_height,
    )


def build_screenshot(slide, prs, image_path: Path, title: str):
    """Slide novo com 1 screenshot fullscreen + badge de título."""
    add_solid_bg(slide, prs, LIGHT_BG)
    slide.shapes.add_picture(
        str(image_path),
        0, 0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    # badge titulo
    badge = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(0.3),
        Inches(3.8), Inches(0.4),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK
    badge.line.fill.background()
    badge.shadow.inherit = False
    tf = badge.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Outfit"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = WHITE


def move_block_to(prs, refs_to_move, dst_idx):
    """Move um conjunto de refs (na ordem) para a posicao dst_idx."""
    xml_slides = prs.slides._sldIdLst
    # remove todas as refs
    for ref in refs_to_move:
        xml_slides.remove(ref)
    # reinsere em ordem
    for offset, ref in enumerate(refs_to_move):
        xml_slides.insert(dst_idx + offset, ref)


def main():
    if not BACKUP.exists():
        shutil.copy(PPT_PATH, BACKUP)
        print(f"  [BKP] backup criado: {BACKUP.name}")
    else:
        print(f"  [BKP] backup ja existe: {BACKUP.name}")

    prs = Presentation(str(PPT_PATH))
    total_inicial = len(prs.slides)
    print(f"  Slides iniciais: {total_inicial}")

    slides_list = list(prs.slides)
    # ─── 1) Transformar os 5 placeholders em divisores (via screenshot HTML) ───
    print("\n  Esvaziando e aplicando screenshots de divisores HTML...")
    for ordem, (nome, prefix, descricao, div_img) in enumerate(TELAS, 1):
        sl_idx = PLACEHOLDER_IDXS[ordem - 1]
        slide = slides_list[sl_idx]
        div_path = DIV_SHOTS / div_img
        if not div_path.exists():
            print(f"     [!] divisor nao encontrado: {div_path}")
            continue
        build_divisor(slide, prs, div_path)
        print(f"     [{ordem}] slide #{sl_idx+1} -> divisor: {nome} ({div_img})")

    # ─── 2) Adicionar 15 prints ao final + capturar refs ───
    print("\n  Adicionando 15 prints ao final...")
    blank_layout = prs.slide_layouts[6]
    novos_por_tela = {}  # prefix -> lista de refs xml dos slides novos

    for ordem, (nome, prefix, descricao, _div) in enumerate(TELAS, 1):
        shot_dir = SHOTS / prefix
        shots = sorted(shot_dir.glob(f"{prefix}-slide*.png"))
        novos_por_tela[prefix] = []
        for i, shot in enumerate(shots, 1):
            slide = prs.slides.add_slide(blank_layout)
            build_screenshot(slide, prs, shot, f"{nome} · parte {i}/{len(shots)}")
            # captura ref XML pra reordenar depois
            xml_slides = prs.slides._sldIdLst
            novos_por_tela[prefix].append(list(xml_slides)[-1])
            print(f"     {nome} · {i}/{len(shots)}: {shot.name}")

    print(f"  Slides após adicao: {len(prs.slides)}")

    # ─── 3) Reordenar: cada print vai logo após seu divisor ───
    # Ordem alvo final:
    #   [0-4 originais] [div Dash, print1, print2, print3, div Prev, print1, print2, print3,
    #    div Casc, p1..p5, div Saude, p1, p2, div Morn, p1, p2]
    #   [10-13 originais: planejamento, recomendacoes]
    #
    # Os divisores estao em pos 5-9 (placeholders esvaziados).
    # Os prints estao em pos 14-28 (no final, na ordem certa).
    # Vou mover cada bloco de prints pra logo apos seu divisor.

    print("\n  Reordenando prints intercalados com divisores...")
    # após cada move, indices mudam. Vou usar approach por referencia:
    # 1. Capturo refs de divisores E prints
    # 2. Construo lista alvo na ordem desejada
    # 3. Remove tudo da sldIdLst, reinsere na ordem correta

    xml_slides = prs.slides._sldIdLst
    all_refs = list(xml_slides)  # ordem atual

    # estrutura desejada (ordem final dos refs):
    # [originais 0..4] + [interleaved divisor+prints] + [originais 10..13]
    originais_inicio = all_refs[0:5]
    divisores = all_refs[5:10]  # indices 5-9, na ordem Dashboard..Morning
    originais_fim = all_refs[10:14]  # planejamento, recomendacoes

    interleaved = []
    for ordem, (nome, prefix, _desc, _div) in enumerate(TELAS):
        interleaved.append(divisores[ordem])
        for print_ref in novos_por_tela[prefix]:
            interleaved.append(print_ref)

    nova_ordem = originais_inicio + interleaved + originais_fim

    # limpa a lista e reinsere
    for ref in list(xml_slides):
        xml_slides.remove(ref)
    for ref in nova_ordem:
        xml_slides.append(ref)

    print(f"  Slides reordenados — total final: {len(prs.slides)}")

    # ─── 4) Salva ───
    prs.save(str(PPT_PATH))
    final_size_mb = PPT_PATH.stat().st_size / 1024 / 1024
    print(f"\n  ✓ PPTX salvo: {PPT_PATH.name}")
    print(f"  Tamanho: {final_size_mb:.1f} MB")
    print(f"\n  Estrutura final:")
    print(f"    Slides 1-5: originais (capa, intro, arquitetura)")
    print(f"    Slides 6-9: Dashboard (divisor + 3 prints)")
    print(f"    Slides 10-13: Previsões (divisor + 3 prints)")
    print(f"    Slides 14-19: Cascatas (divisor + 5 prints)")
    print(f"    Slides 20-22: Saúde · Produto (divisor + 2 prints)")
    print(f"    Slides 23-25: Morning brief (divisor + 2 prints)")
    print(f"    Slides 26-29: originais (planejamento, recomendações)")


if __name__ == "__main__":
    main()
