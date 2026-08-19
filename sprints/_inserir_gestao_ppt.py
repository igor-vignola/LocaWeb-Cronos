"""
Insere os 3 screenshots da secao Gestao no PPTX.

Estrategia: esvazia os 3 placeholders da secao "4. PLANEJAMENTO E GESTAO" (slides 21-23)
e aplica os screenshots como imagem fullscreen.

Uso:
    python sprints/_inserir_gestao_ppt.py
"""
import sys
from pathlib import Path

from pptx import Presentation

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

SPRINTS = Path(__file__).parent
PROTO = SPRINTS.parent / "prototipos"
PPT_PATH = SPRINTS / "EC_Sprint_2_2TSCOA_arqsolucao_Cronos_SuperDataBros.pptx"
GESTAO_SHOTS = PROTO / "gestao" / "screenshots"

# (indice_0based_do_placeholder, arquivo_png)
SUBSTITUICOES = [
    (20, "intro.png"),    # slide 21 — "4. PLANEJAMENTO E GESTAO DE PROJETOS"
    (21, "board.png"),    # slide 22 — "4. PLANEJAMENTO E GESTAO DE PROJETOS" (Exemplo)
    (22, "card.png"),     # slide 23 — (sem texto)
]


def clear_slide(slide):
    """Remove TODAS as shapes do slide."""
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        tag = sp.tag.split("}")[-1]
        if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
            spTree.remove(sp)


def main():
    prs = Presentation(str(PPT_PATH))
    total = len(prs.slides)
    print(f"  Slides totais no PPT: {total}")

    slides_list = list(prs.slides)
    print(f"\n  Aplicando 3 screenshots da secao Gestao...")

    for idx, png_name in SUBSTITUICOES:
        png_path = GESTAO_SHOTS / png_name
        if not png_path.exists():
            print(f"     [ERR] arquivo nao encontrado: {png_path}")
            continue
        if idx >= total:
            print(f"     [ERR] slide #{idx+1} fora do range (total={total})")
            continue

        slide = slides_list[idx]
        clear_slide(slide)
        slide.shapes.add_picture(
            str(png_path),
            0, 0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"     slide #{idx+1} -> {png_name}")

    prs.save(str(PPT_PATH))
    final_mb = PPT_PATH.stat().st_size / 1024 / 1024
    print(f"\n  PPTX salvo: {PPT_PATH.name}")
    print(f"  Tamanho: {final_mb:.1f} MB")


if __name__ == "__main__":
    main()
