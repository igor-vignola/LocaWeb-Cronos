"""
Inspeciona o PPTX da Sprint 2 — lista slides, layouts, dimensoes e textos.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

PPT = Path(__file__).parent / "EC_Sprint_2_2TSCOA_arqsolucao_Cronos_SuperDataBros.pptx"

prs = Presentation(str(PPT))

w_in = prs.slide_width / 914400
h_in = prs.slide_height / 914400
print(f"Slide size: {w_in:.2f} x {h_in:.2f} in  ({prs.slide_width} x {prs.slide_height} EMU)")
print(f"Total slides: {len(prs.slides)}")
print()

# layouts disponiveis no master
print("=== Layouts disponiveis ===")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  [{i}] {layout.name}")
print()

# listar todos os slides com seus textos
print("=== Slides do PPT atual ===")
for idx, slide in enumerate(prs.slides, 1):
    layout = slide.slide_layout.name
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs).strip()
                if txt:
                    texts.append(txt)
    title = texts[0] if texts else "(sem texto)"
    snippet = " / ".join(texts[1:4])[:90] if len(texts) > 1 else ""
    print(f"  #{idx:2d} [layout: {layout:<25}] {title[:60]}")
    if snippet:
        print(f"      {snippet}")
