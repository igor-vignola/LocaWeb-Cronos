"""
Insere o slide 'Demo navegavel online' no PPTX da Sprint 2.

Adiciona como novo slide no final e move pra posicao 24 (antes do
'RECOMENDACOES FINAIS').

Uso:
    python sprints/_inserir_demo_ppt.py
"""
import sys
from pathlib import Path

from pptx import Presentation

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

SPRINTS = Path(__file__).parent
PROTO = SPRINTS.parent / "prototipos"
PPT_PATH = SPRINTS / "EC_Sprint_2_2TSCOA_arqsolucao_Cronos_SuperDataBros.pptx"
DEMO_PNG = PROTO / "gestao" / "screenshots" / "demo.png"

# Insere ANTES do slide RECOMENDACOES FINAIS (que era o ultimo, indice = total - 1).
# Apos add_slide e antes do move, posicao alvo eh "total - 1" (empurra
# recomendacoes pra ultimo).
INSERT_BEFORE_LAST = True


def main():
    prs = Presentation(str(PPT_PATH))
    total = len(prs.slides)
    print(f"  Slides totais antes: {total}")

    # adiciona slide ao final
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(DEMO_PNG),
        0, 0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    print(f"  Slide demo adicionado (atual ultimo, indice {total})")

    # move pra antes do ultimo (que era Recomendacoes)
    if INSERT_BEFORE_LAST:
        xml_slides = prs.slides._sldIdLst
        all_refs = list(xml_slides)
        new_ref = all_refs[-1]   # o que acabei de adicionar
        target_pos = total - 1   # posicao do antigo ultimo (RECOMENDACOES)
        xml_slides.remove(new_ref)
        xml_slides.insert(target_pos, new_ref)
        print(f"  Movido para posicao {target_pos + 1} (antes do RECOMENDACOES FINAIS)")

    prs.save(str(PPT_PATH))
    final_mb = PPT_PATH.stat().st_size / 1024 / 1024
    print(f"\n  PPTX salvo: {PPT_PATH.name}")
    print(f"  Slides totais agora: {len(prs.slides)}")
    print(f"  Tamanho: {final_mb:.1f} MB")


if __name__ == "__main__":
    main()
